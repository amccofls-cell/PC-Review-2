# -*- coding: utf-8 -*-
"""
PPTX 비교표 추출 모듈 — python-pptx 로 슬라이드 내 모든 표를 순회하며
행/열/셀/병합 정보를 보존해 추출한다(명세서 6장 STEP 3, 8장).

반환 형식(모든 파서 공통의 "테이블" dict):
{
  "source": "pptx",
  "slide": 1,                # 1-based 슬라이드 번호
  "table_index": 0,          # 슬라이드 내 표 순번
  "rows": [ [ {"text","rowspan","colspan"}, ... ], ... ]
}
"""
import io

from pptx import Presentation
from pptx.util import Emu


class PptxParseError(Exception):
    """PPTX 읽기/파싱 오류 (사용자에게 그대로 노출)"""


def _cell_text(cell):
    """셀 내 모든 텍스트(병합 포함)를 줄바꿈으로 이어붙인다."""
    if cell is None:
        return ""
    texts = []
    for p in cell.text_frame.paragraphs:
        t = "".join(run.text for run in p.runs).strip()
        if t:
            texts.append(t)
    return "\n".join(texts)


def _read_table(table):
    rows = []
    nrows, ncols = len(table.rows), len(table.columns)
    for r in range(nrows):
        row = []
        for c in range(ncols):
            cell = table.cell(r, c)
            rowspan, colspan = 1, 1
            # python-pptx 병합 셀: 병합 시작점(is_merge_origin)에만 span 크기 정보
            if cell.is_merge_origin:
                colspan = cell.span_width
                rowspan = cell.span_height
            row.append({
                "text": _cell_text(cell),
                "rowspan": rowspan,
                "colspan": colspan,
            })
        rows.append(row)
    return rows


def parse_pptx(file_bytes):
    """
    PPTX 파일 바이트를 받아 모든 슬라이드의 모든 표를 추출한다.
    file_bytes: bytes (업로드 파일의 getvalue())
    """
    if not file_bytes:
        raise PptxParseError("빈 파일입니다. PPTX 파일을 선택해 주세요.")
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        raise PptxParseError(f"PPTX 파일을 열지 못했습니다. 형식을 확인해 주세요. ({exc})")
    tables = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        ti = 0
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            try:
                rows = _read_table(shape.table)
            except Exception as exc:
                raise PptxParseError(f"{slide_idx}번 슬라이드의 표를 읽지 못했습니다. ({exc})")
            tables.append({
                "source": "pptx",
                "slide": slide_idx,
                "table_index": ti,
                "rows": rows,
            })
            ti += 1
    if not tables:
        raise PptxParseError("PPTX 파일에서 표를 찾지 못했습니다. 표가 포함된 슬라이드를 확인해 주세요.")
    return tables
