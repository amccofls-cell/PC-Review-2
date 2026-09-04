# -*- coding: utf-8 -*-
"""
오프라인 자가 테스트 — API 키 없이 검증 가능한 항목을 모두 테스트한다.

실제 MFDS/HIRA 호출(명세서 14장 테스트 1·2)은 API 키가 필요하므로 여기서 실행하지 않고
스킵 목록에 표시한다. 나머지(3~11, 13·14 일부)는 실제 코드 경로로 검증한다.

실행:  python tests/selftest.py
"""
import io
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from openpyxl import Workbook
from pptx import Presentation

from modules import (
    mfds_api, hira_api, drug_matcher, pptx_parser, xlsx_parser, clipboard_parser,
    table_normalizer as normalizer, rule_validator, claude_prompt_builder, result_parser,
    cache_store, grouping,
)

PASS, FAIL = "PASS", "FAIL"
results = []


def check(name, cond, detail=""):
    results.append((name, PASS if cond else FAIL, detail))
    print(f"[{'✓' if cond else '✗'}] {name}" + (f" — {detail}" if detail else ""))


def make_sample_pptx():
    """표 2개(일반 + 병합 셀)를 포함한 샘플 PPTX 생성."""
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 3, 3
    shape = slide.shapes.add_table(rows, cols, 10, 10, 600, 120)
    tbl = shape.table
    data = [
        ["항목", "신청의약품", "비교의약품"],
        ["제품명", "온코정 100mg", "온코정 50mg"],
        ["용법용량", "1일 1회 500mg", "1일 1회 250mg"],
    ]
    for r in range(rows):
        for c in range(cols):
            tbl.cell(r, c).text = data[r][c]
    # 병합 셀 포함 표 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    shape2 = slide2.shapes.add_table(2, 2, 10, 10, 400, 80)
    t2 = shape2.table
    t2.cell(0, 0).text = "병합시작"
    t2.cell(0, 1).text = "B1"
    t2.cell(1, 0).text = "A2"
    prs.save(buf)
    return buf.getvalue()


def make_sample_xlsx():
    """병합 셀이 포함된 샘플 XLSX 생성."""
    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "비교표"
    ws.append(["항목", "신청의약품", "비교의약품"])
    ws.append(["효능효과", "전이성 위암 치료", "전이성 위암 2차 치료"])
    ws.append(["이상반응", "오심, 구토", "오심"])
    ws.merge_cells("A3:A4")
    ws["A3"] = "병합된 항목"
    wb.save(buf)
    return buf.getvalue()


def main():
    print("=" * 60)
    print("의약품 심의자료 검증기 v4.0 — 오프라인 자가 테스트")
    print("=" * 60)

    # ---- 1) PPTX 파서 ----
    pptx_bytes = make_sample_pptx()
    tables = pptx_parser.parse_pptx(pptx_bytes)
    check("3. PPTX 표 추출 (슬라이드 순회)", len(tables) == 2, f"{len(tables)}개 표")
    check("3. PPTX 행/열 구조", tables[0]["rows"][0][1]["text"] == "신청의약품"
          and tables[0]["rows"][2][1]["text"] == "1일 1회 500mg")

    # ---- 2) XLSX 파서 + 병합 ----
    xlsx_bytes = make_sample_xlsx()
    xt = xlsx_parser.parse_xlsx(xlsx_bytes)[0]
    check("4. XLSX 표 추출 (시트)", xt["slide"] == "비교표" and len(xt["rows"]) == 4)
    merged_ok = any(c["rowspan"] >= 2 for row in xt["rows"] for c in row if c["text"] == "병합된 항목")
    check("4. XLSX 병합 셀 보존", merged_ok)

    # ---- 3) 복사·붙여넣기 파서 ----
    pasted = "항목\t신청의약품\t비교의약품\n제품명\t온코정 100mg\t온코정 50mg\n용법용량\t1일 1회 500mg\t1일 1회 250mg"
    pt = clipboard_parser.parse_clipboard(pasted)[0]
    check("5. 붙여넣기 표 인식", len(pt["rows"]) == 3 and pt["rows"][0][2]["text"] == "비교의약품")
    try:
        clipboard_parser.parse_clipboard("")
        check("5. 빈 붙여넣기 오류 처리", False)
    except clipboard_parser.ClipboardParseError:
        check("5. 빈 붙여넣기 오류 처리", True)

    # ---- 4) 공통 스키마 정규화 + 방향 추정 ----
    orient = normalizer.guess_orientation(xt["rows"])
    check("4. 방향 자동 추정 (행=항목)", orient == normalizer.ORIENT_ROWS_ARE_ITEMS)
    pairs = normalizer.to_field_product_pairs(xt, orient)
    schema_ok = all(set(p) >= {"slide", "table_index", "row", "column", "field", "product", "value"} for p in pairs)
    check("8. 공통 스키마(명세서 8장) 생성", schema_ok and len(pairs) >= 2, f"{len(pairs)}개 셀")
    field_ok = any(p["field"] == "효능효과" for p in pairs)
    check("8. 필드 키워드 인식(효능효과)", field_ok)

    # ---- 5) 제품 매칭 (8자리 바코드 규칙) ----
    detail = {"바코드(표준코드)": "8801234500128", "제품명": "온코정 100mg"}
    hira_rows = [
        {"mdsCd": "1234500101234", "itmNm": "온코정 100mg", "mnfEntpNm": "온코팜", "mxCprc": "50000"},
        {"mdsCd": "9999999900000", "itmNm": "다른약", "mnfEntpNm": "X", "mxCprc": "100"},
    ]
    assert drug_matcher.barcode_8digits("8801234500128") == "12345001"
    match = drug_matcher.match_hira(detail, hira_rows)
    check("2. 바코드 8자리 매칭(1순위)", match["status"] == drug_matcher.STATUS_BARCODE, match["status"])
    partial = drug_matcher.match_hira({"바코드(표준코드)": "", "제품명": "온코정"}, hira_rows)
    check("2. 품목명 부분 일치 → 🟠 확인 필요", partial["status"] == drug_matcher.STATUS_NAME_PARTIAL, partial["status"])
    fail = drug_matcher.match_hira({"바코드(표준코드)": "", "제품명": "황산마그네슘주 10%"}, hira_rows)
    check("2. 매칭 실패 시 자동선택 금지", fail["status"] == drug_matcher.STATUS_FAIL, fail["status"])
    b8 = drug_matcher.barcode_8digits("8801234500128")
    check("2. 바코드 슬라이싱 [3:11]", b8 == "12345001", f"b8={b8}")

    # ---- 6) MFDS 모듈 단위 (XML 파싱, 섹션 분리) ----
    ee_xml = "<DOC><TITLE>효능효과</TITLE><TEXT>위암 환자의 치료</TEXT></DOC>"
    ee = mfds_api.docs_to_text(mfds_api.parse_doc_xml(ee_xml))
    check("1. MFDS EE_DOC_DATA XML 파싱", "위암 환자의 치료" in ee)
    nb_xml = (
        "<DOC>"
        "<TITLE>1. 다음 환자에게는 투여하지 말 것</TITLE><TEXT>이 약에 과민성인 환자</TEXT>"
        "<TITLE>2. 이상반응</TITLE><TEXT>오심, 구토</TEXT>"
        "<TITLE>(공고번호)3. 소아에 대한 투여</TITLE><TEXT>안전성 미확립</TEXT>"
        "</DOC>"
    )
    nb = mfds_api.split_nb_sections(mfds_api.parse_doc_xml(nb_xml))
    check("1. NB_DOC_DATA 섹션 분리(금기)", "과민성인 환자" in nb.get("금기사항", ""))
    check("1. NB_DOC_DATA 섹션 분리(이상반응)", "오심, 구토" in nb.get("이상반응", ""))
    check("1. NB 타이틀 정규화(번호·괄호 제거)", "소아에 대한 투여" in nb.get("소아_고령자투여", ""))
    bad_xml = "<DOC><TITLE>효능효과</TITLE><TEXT>닫힘 태그 누락"
    raw = mfds_api.docs_to_text(mfds_api.parse_doc_xml(bad_xml))
    check("1. 잘못된 XML → 원문 유지", "효능효과" in raw or "닫힘 태그" in raw)

    # ---- 7) HIRA 모듈 단위 ----
    check("2. HIRA 삭제 행 필터(_filter_rows는 응답 봉투 필요, 행 구성만 확인)",
          hira_api.get_price({"mxCprc": "12,345"}) == 12345.0)
    check("2. HIRA 가격 파싱(빈 값)", hira_api.get_price({"mxCprc": ""}) is None)

    # ---- 8) 규칙 검증 ----
    st1, r1 = rule_validator.check_basic_info("온코정 100mg", "온코정 100mg")
    check("6. 기본정보 일치", st1 == rule_validator.STATUS_OK, r1)
    st2, r2 = rule_validator.check_basic_info("(주)온코팜", "주식회사 온코팜")
    check("6. 법인 표기 차이 무시", st2 == rule_validator.STATUS_OK, r2)
    st3, r3 = rule_validator.check_basic_info("온코정 100mg", "온코정 50mg")
    check("6. 기본정보 불일치 검출", st3 == rule_validator.STATUS_FIX, r3)

    st4, r4 = rule_validator.check_numeric_field("1일 1회 500mg", "1일 1회 5000mg")
    check("6. 숫자·단위 불일치 검출", st4 == rule_validator.STATUS_FIX, r4)
    st5, r5 = rule_validator.check_numeric_field("1일 1회 500mg", "1일 1회 500mg")
    check("6. 숫자 일치 → Claude 확인 필요로 전환", st5 == rule_validator.STATUS_CLAUDE, r5)

    st6, r6, _ = rule_validator.check_price("50,000", 50000.0)
    check("7. 약가 일치 판정", st6 == rule_validator.STATUS_OK, r6)
    st7, r7, _ = rule_validator.check_price("55,000", 50000.0)
    check("7. 약가 불일치 검출", st7 == rule_validator.STATUS_FIX, r7)
    diff = rule_validator.price_diff_percent(55000, 50000)
    check("7. 약가 차이율 공식 (55000-50000)/50000*100", diff == 10.0, f"diff={diff}")

    # ---- 9) Claude 프롬프트 빌더 ----
    fake_detail = {"제품명": "온코정", "바코드(표준코드)": "8801234500128", "효능효과": "위암 환자의 치료",
                   "용법용량": "1일 1회 500mg", "성분명": "온코정성분", "제조판매사": "온코팜",
                   "사용상주의사항": {"이상반응": "오심, 구토"}, "error": None}
    prompt_prods = [
        {"label": "신청의약품 | 온코팜 | 1234", "detail": fake_detail, "hira_rows": [
            {"itmNm": "온코정", "mnfEntpNm": "온코팜", "mdsCd": "1234500100000", "mxCprc": "50000", "payTpNm": "급여", "meftDivNo": "421"}]},
    ]
    full = claude_prompt_builder.build_full_prompt(prompt_prods)
    check("8. 프롬프트 구조([역할]/[검증 원칙]/[검증 요청]/[출력 형식])",
          all(k in full for k in ("[역할]", "[검증 원칙]", "[검증 요청]", "[출력 형식]")),
          f"{len(full)}자")
    check("8. MFDS 원문 미요약 포함", "위암 환자의 치료" in full and "1일 1회 500mg" in full)
    check("8. HIRA 약가정보 포함", "상한금액: 50000원" in full)
    field_p = claude_prompt_builder.build_field_prompt("이상반응", prompt_prods)
    check("8. 항목별 프롬프트(이상반응 원문만)", "오심, 구토" in field_p)
    field_p2 = claude_prompt_builder.build_field_prompt("효능효과", prompt_prods)
    check("8. 항목별 프롬프트(효능효과)", "위암 환자의 치료" in field_p2)

    # ---- 10) Claude 결과 재붙여넣기 파싱 ----
    md = "| 제품 | 항목 | 판단 | 이유 | 원문 근거 |\n|---|---|---|---|---|\n| 온코정 | 효능효과 | 수정필요 | 대상자 조건 누락 | 원문: 위암 환자의 치료 |"
    pr = result_parser.parse_result(md)
    check("11. 마크다운 결과 파싱", pr["ok"] and pr["rows"][0]["판단"] == "수정필요", str(pr.get("rows", []))[:60])
    js = '{"results": [{"제품": "온코정", "항목": "용법용량", "판단": "일치", "이유": "동일", "원문": "1일 1회 500mg"}]}'
    pj = result_parser.parse_result(js)
    check("11. JSON 결과 파싱", pj["ok"] and pj["rows"][0]["원문근거"] == "1일 1회 500mg")
    pbad = result_parser.parse_result("이건 표도 JSON도 아닙니다")
    check("11. 파싱 실패 메시지 (추측 금지)", (not pbad["ok"]) and "형식을 인식하지 못했습니다" in pbad["error"])

    # ---- 11) Claude API 호출 코드 부재 (명세서 14-13) ----
    banned = []
    for root, _, files in os.walk(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))):
        if "__pycache__" in root or root.endswith(".git"):
            continue
        for fn in files:
            if not fn.endswith(".py"):
                continue
            path = os.path.join(root, fn)
            with open(path, encoding="utf-8") as f:
                content = f.read().lower()
                if ("anthro" + "pic") in content or ("api.anthro" + "pic.com") in content:
                    banned.append(path)
    check("13. 금지 API 호출 코드 0건 (grep 검사)", len(banned) == 0,
          banned if banned else "전체 파이썬 소스에서 0건")

    # ---- 12) API 키 하드코딩 부재 (명세서 14-14) ----
    hardcoded = []
    for root, _, files in os.walk(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))):
        if "__pycache__" in root or root.endswith(".git"):
            continue
        for fn in files:
            if not fn.endswith(".py"):
                continue
            path = os.path.join(root, fn)
            with open(path, encoding="utf-8") as f:
                for ln in f:
                    if ("serviceKey" in ln and "st.secrets" not in ln and "params" not in ln
                            and "param" not in ln and "serviceKey\"" not in ln):
                        # 실제 키처럼 보이는 값(긴 알파벳숫자/하이픈)만 플래그
                        if re_search_key(ln):
                            hardcoded.append((path, ln.strip()[:80]))
    check("14. API 키 하드코딩 없음", len(hardcoded) == 0, hardcoded if hardcoded else "소스에 키 값 없음")

    # ---- 15) v4.1 캐시 저장/로드 + 로컬 필터 (API 호출 없음) ----
    cache_items = [
        {"ITEM_SEQ": "A1", "ITEM_NAME": "온코정 100mg", "ENTP_NAME": "온코팜", "ITEM_PERMIT_DATE": "2020-01-01"},
        {"ITEM_SEQ": "A2", "ITEM_NAME": "포비정", "ENTP_NAME": "한독약품", "ITEM_PERMIT_DATE": "2021-02-02"},
    ]
    flt = cache_store.filter_products(cache_items, "온코")
    check("15. 로컬 키워드 필터(제품명)", len(flt) == 1 and flt[0]["ITEM_SEQ"] == "A1")
    flt2 = cache_store.filter_products(cache_items, "한독")
    check("15. 로컬 키워드 필터(제조사명)", len(flt2) == 1 and flt2[0]["ITEM_SEQ"] == "A2")
    check("15. 빈 키워드 → 전체 반환", len(cache_store.filter_products(cache_items, "")) == 2)
    tmp_cache = os.path.join(tempfile.gettempdir(), "selftest_cache_v41.json")
    saved = cache_store.save_cache(tmp_cache, cache_items)
    loaded = cache_store.load_cache(tmp_cache)
    check("15. JSON 캐시 저장/로드 왕복", loaded is not None and loaded["count"] == 2
          and loaded["items"][0]["ITEM_NAME"] == "온코정 100mg" and saved["count"] == 2)

    # ---- 16) v4.1 mock 응답 — 전체 적재(품목/약가) + 상세 파싱 ----
    from unittest import mock as _umock

    def _make_resp(payload):
        class _R:
            def __init__(self, d):
                self._d = d
            def json(self):
                return self._d
        return _R(payload)

    mfds_payload = {
        "header": {"resultCode": "00", "resultMsg": "NORMAL"},
        "body": {"totalCount": 2, "items": {"item": [
            {"ITEM_SEQ": "1001", "ITEM_NAME": "온코정 100mg", "ENTP_NAME": "온코팜",
             "CANCEL_NAME": "정상", "ITEM_PERMIT_DATE": "2020-01-01"},
            {"ITEM_SEQ": "1002", "ITEM_NAME": "폐기약", "ENTP_NAME": "X",
             "CANCEL_NAME": "취하", "ITEM_PERMIT_DATE": ""},
        ]}},
    }
    with _umock.patch("modules.mfds_api.requests.get", return_value=_make_resp(mfds_payload)):
        fetched = mfds_api.fetch_all_products("TESTKEY")
    check("16. MFDS 전체 적재(mock) — 정상품목만", len(fetched) == 1 and fetched[0]["ITEM_SEQ"] == "1001")

    hira_payload = {
        "header": {"resultCode": "00", "resultMsg": "NORMAL"},
        "body": {"totalCount": 2, "items": {"item": [
            {"mdsCd": "1234500100000", "itmNm": "온코정 100mg", "mnfEntpNm": "온코팜",
             "mxCprc": "50000", "payTpNm": "급여", "meftDivNo": "421"},
            {"mdsCd": "9999900000000", "itmNm": "삭제약", "mnfEntpNm": "X",
             "mxCprc": "1", "payTpNm": "삭제", "meftDivNo": ""},
        ]}},
    }
    with _umock.patch("modules.hira_api.requests.get", return_value=_make_resp(hira_payload)):
        hp = hira_api.fetch_all_drug_prices("TESTKEY")
    check("16. HIRA 전체 적재(mock) — 삭제 행 제외", len(hp) == 1 and hp[0]["mdsCd"] == "1234500100000")

    detail_payload = {
        "header": {"resultCode": "00"},
        "body": {"items": [
            {"ITEM_SEQ": "1001", "ITEM_NAME": "온코정 100mg", "ENTP_NAME": "온코팜",
             "MAIN_ITEM_INGR": "[A12345]온코정성분", "BAR_CODE": "8801234500128",
             "EE_DOC_DATA": "<DOC><TITLE>효능효과</TITLE><TEXT>위암 환자의 치료</TEXT></DOC>",
             "UD_DOC_DATA": "", "NB_DOC_DATA": "", "ITEM_ENG_NAME": "", "MAIN_INGR_ENG": "",
             "ETC_OTC_CODE": "", "ATC_CODE": "", "MATERIAL_NAME": "", "PACK_UNIT": "",
             "VALID_TERM": "", "STORAGE_METHOD": "", "CHART": "", "EDI_CODE": ""},
        ]},
    }
    with _umock.patch("modules.mfds_api.requests.get", return_value=_make_resp(detail_payload)):
        det = mfds_api.get_product_detail("1001", "TESTKEY")
    check("16. MFDS 상세(mock) — 성분명 대괄호 코드 제거", det["성분명"] == "온코정성분")
    check("16. MFDS 상세(mock) — 효능효과 XML 파싱(원문 유지)", "위암 환자의 치료" in det["효능효과"])
    check("16. MFDS 상세(mock) — 바코드 보존", det["바코드(표준코드)"] == "8801234500128")

    # ---- 17) 그룹 데이터모델 + 4그룹/3강도 시나리오 ----
    legacy_groups = grouping.migrate_legacy(selection=["A100", "A200", "C100"], applicant_seq="A100", comparator_seqs=["A200"])
    check("17. 구버전 상태 마이그레이션", legacy_groups["applicant"]["seqs"] == ["A100"] and "C100" in legacy_groups["comp1"]["seqs"])

    groups = grouping.ensure_minimum_groups({
        "applicant": {"label": "신청의약품", "seqs": ["N100", "N200", "N300"]},
        "comp1": {"label": "비교의약품1", "seqs": ["F100", "F200", "F400"]},
        "comp2": {"label": "비교의약품2", "seqs": ["A100", "A200", "A400"]},
        "comp3": {"label": "비교의약품3", "seqs": ["I050", "I100", "I200"]},
    })
    by_seq = {
        "N100": {"ITEM_SEQ": "N100", "ITEM_NAME": "나르코설하정100마이크로그램", "ENTP_NAME": "신청사"},
        "N200": {"ITEM_SEQ": "N200", "ITEM_NAME": "나르코설하정200마이크로그램", "ENTP_NAME": "신청사"},
        "N300": {"ITEM_SEQ": "N300", "ITEM_NAME": "나르코설하정300마이크로그램", "ENTP_NAME": "신청사"},
        "F100": {"ITEM_SEQ": "F100", "ITEM_NAME": "펜토라박칼정100마이크로그램", "ENTP_NAME": "비교사1"},
        "F200": {"ITEM_SEQ": "F200", "ITEM_NAME": "펜토라박칼정200마이크로그램", "ENTP_NAME": "비교사1"},
        "F400": {"ITEM_SEQ": "F400", "ITEM_NAME": "펜토라박칼정400마이크로그램", "ENTP_NAME": "비교사1"},
        "A100": {"ITEM_SEQ": "A100", "ITEM_NAME": "앱스트랄설하정100마이크로그램", "ENTP_NAME": "비교사2"},
        "A200": {"ITEM_SEQ": "A200", "ITEM_NAME": "앱스트랄설하정200마이크로그램", "ENTP_NAME": "비교사2"},
        "A400": {"ITEM_SEQ": "A400", "ITEM_NAME": "앱스트랄설하정400마이크로그램", "ENTP_NAME": "비교사2"},
        "I050": {"ITEM_SEQ": "I050", "ITEM_NAME": "인스타닐나잘스프레이50마이크로그램", "ENTP_NAME": "비교사3"},
        "I100": {"ITEM_SEQ": "I100", "ITEM_NAME": "인스타닐나잘스프레이100마이크로그램", "ENTP_NAME": "비교사3"},
        "I200": {"ITEM_SEQ": "I200", "ITEM_NAME": "인스타닐나잘스프레이200마이크로그램", "ENTP_NAME": "비교사3"},
    }
    group_items = grouping.build_group_items(groups, by_seq)
    check("17. 4그룹/3강도 시나리오 품목 수", len(group_items) == 12, f"{len(group_items)}건")
    strengths = {row["seq"]: row["strength"] for row in group_items}
    check("17. 제품명 기반 함량 파싱", strengths["N100"] == "100mcg" and strengths["F400"] == "400mcg" and strengths["I050"] == "50mcg", str(strengths))
    reordered = grouping.assign_seqs_to_group(groups, "comp3", ["N300"])
    check("17. 그룹 재지정 시 중복 제거", "N300" not in reordered["applicant"]["seqs"] and reordered["comp3"]["seqs"][-1] == "N300")

    prompt_prods_group = [
        {"label": "X", "prompt_label": "신청의약품 / 100mcg / 나르코설하정100마이크로그램", "role": "신청의약품", "detail": fake_detail, "hira_rows": [], "pairs": []},
        {"label": "Y", "prompt_label": "비교의약품1 / 100mcg / 펜토라박칼정100마이크로그램", "role": "비교의약품1", "detail": fake_detail, "hira_rows": [], "pairs": []},
        {"label": "Z", "prompt_label": "비교의약품2 / 100mcg / 앱스트랄설하정100마이크로그램", "role": "비교의약품2", "detail": fake_detail, "hira_rows": [], "pairs": []},
        {"label": "W", "prompt_label": "비교의약품3 / 50mcg / 인스타닐나잘스프레이50마이크로그램", "role": "비교의약품3", "detail": fake_detail, "hira_rows": [], "pairs": []},
    ]
    full_group = claude_prompt_builder.build_full_prompt(prompt_prods_group)
    check("17. 그룹 프롬프트 대상열거", all(x in full_group for x in ("신청의약품:", "비교의약품1:", "비교의약품2:", "비교의약품3:")))

    # ---- 18) v4.1+ 오류 봉투·페이지네이션·키 인코딩 ----
    key_enc = mfds_api.prep_service_key("abc%2Bdef%3Dg")
    check("18. serviceKey 이중 인코딩 방지(unquote 1회)", key_enc == "abc+def=g", key_enc)
    err_payload = {"OpenAPI_ServiceResponse": {"cmmMsgHeader": {
        "errMsg": "SERVICE_KEY_IS_NOT_REGISTERED_ERROR",
        "returnAuthMsg": "등록되지 않은 서비스키", "returnReasonCode": "30"}}}
    try:
        with _umock.patch("modules.mfds_api.requests.get", return_value=_make_resp(err_payload)):
            mfds_api.fetch_all_products("BADKEY")
        check("18. 오류 봉투 → 명확한 한국어 메시지", False)
    except mfds_api.MfdsApiError as e:
        check("18. 오류 봉투 → 명확한 한국어 메시지", "등록되지 않은 서비스키" in str(e), str(e))

    # 페이지네이션: totalCount=1500 / numOfRows=500 → 3페이지, mock 호출 3번
    pages = []
    for pn in (1, 2, 3):
        pages.append(_make_resp({"header": {"resultCode": "00"}, "body": {"totalCount": 1500, "items": {"item": [
            {"ITEM_SEQ": f"S{pn}{i}", "ITEM_NAME": f"품목{pn}-{i}", "ENTP_NAME": "온코팜", "CANCEL_NAME": "정상", "ITEM_PERMIT_DATE": ""}
            for i in range(500)]}}}))
    call_count = {"n": 0}

    def _side(*a, **k):
        call_count["n"] += 1
        return pages[call_count["n"] - 1]

    with _umock.patch("modules.mfds_api.requests.get", side_effect=_side):
        pag = mfds_api.fetch_all_products("K")
    check("18. 페이지네이션 3페이지 수집", len(pag) == 1500 and call_count["n"] == 3, f"{len(pag)}건/{call_count['n']}회")

    print()
    print("=" * 60)
    n_pass = sum(1 for _, s, _ in results if s == PASS)
    n_fail = sum(1 for _, s, _ in results if s == FAIL)
    print(f"결과: {n_pass} PASS / {n_fail} FAIL / 총 {len(results)}건")
    if n_fail:
        print("실패 항목:")
        for name, s, d in results:
            if s == FAIL:
                print(f"  ✗ {name} — {d}")
    print("=" * 60)
    print("API 키 의존 테스트(명세서 14장의 1·2번)는 실제 키 없이 실행하지 않았습니다 — 배포 후 확인 필요.")
    sys.exit(1 if n_fail else 0)


def re_search_key(line):
    import re
    m = re.search(r"[\"']([A-Za-z0-9+/=_-]{20,})[\"']", line)
    return bool(m)


if __name__ == "__main__":
    main()
